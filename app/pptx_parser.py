"""PPTX and PDF parsing module for CourseAlign API."""
from pptx import Presentation
import fitz  # PyMuPDF
from typing import List, Dict, Any
import os


class PPTXParser:
    """Handles PPTX and PDF text extraction for slides."""
    
    def extract_text_from_slides(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from slides file (PPTX or PDF).
        
        Args:
            file_path: Path to PPTX or PDF file
            
        Returns:
            Dict with slide_texts (list of dicts) and full_text (concatenated)
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    def extract_text_from_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract text from PPTX file including titles, bullets, and speaker notes.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Dict with slide_texts (list of dicts) and full_text (concatenated)
        """
        prs = Presentation(pptx_path)
        slides_data = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            
            # Extract title
            if slide.shapes.title:
                title = slide.shapes.title.text
                if title.strip():
                    slide_text_parts.append(f"Title: {title}")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Skip if it's the title (already extracted)
                    if shape == slide.shapes.title:
                        continue
                    slide_text_parts.append(shape.text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_text_parts.append(f"Speaker Notes: {notes_text}")
            
            slide_text = "\n".join(slide_text_parts)
            
            slides_data.append({
                "slide_number": slide_num,
                "text": slide_text
            })
        
        # Concatenate all slide text
        full_text = "\n\n---SLIDE BREAK---\n\n".join(
            [f"Slide {s['slide_number']}:\n{s['text']}" for s in slides_data]
        )
        
        return {
            "slides": slides_data,
            "full_text": full_text,
            "slide_count": len(slides_data)
        }
    
    def extract_text_from_pdf(self, pdf_path: str) -> Dict[str, Any]:
        """
        Extract text from PDF slides (one page = one slide).
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            Dict with slides (list of dicts) and full_text (concatenated)
        """
        doc = fitz.open(pdf_path)
        slides_data = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            slides_data.append({
                "slide_number": page_num + 1,
                "text": text
            })
        
        doc.close()
        
        # Concatenate all slide text
        full_text = "\n\n---SLIDE BREAK---\n\n".join(
            [f"Slide {s['slide_number']}:\n{s['text']}" for s in slides_data]
        )
        
        return {
            "slides": slides_data,
            "full_text": full_text,
            "slide_count": len(slides_data)
        }
    
    def extract_key_concepts(self, slide_text: str) -> List[str]:
        """
        Extract key concepts from slide text.
        This is a simple extraction based on headings and bullet points.
        
        Args:
            slide_text: Full text from slides
            
        Returns:
            List of key concept strings
        """
        concepts = []
        lines = slide_text.split("\n")
        
        for line in lines:
            line = line.strip()
            # Extract titles
            if line.startswith("Title:"):
                concept = line.replace("Title:", "").strip()
                if concept and len(concept) > 3:
                    concepts.append(concept)
            # Extract bullet points (typically start with -, •, or *)
            elif line and (line[0] in ['-', '•', '*'] or line.startswith("- ")):
                concept = line.lstrip('-•* ').strip()
                if concept and len(concept) > 3:
                    concepts.append(concept)
        
        # Remove duplicates while preserving order
        seen = set()
        unique_concepts = []
        for concept in concepts:
            concept_lower = concept.lower()
            if concept_lower not in seen:
                seen.add(concept_lower)
                unique_concepts.append(concept)
        
        return unique_concepts
